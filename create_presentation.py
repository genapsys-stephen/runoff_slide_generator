from __future__ import print_function
from pptx import Presentation
import pickle
import os.path
from googleapiclient.discovery import build
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from apiclient.http import MediaFileUpload
import json
import platform
import os

file = open('./config.json')
config = json.load(file)

SCOPES = 'https://www.googleapis.com/auth/presentations', 'https://www.googleapis.com/auth/drive', 'https://www.googleapis.com/auth/drive.file',

drive_save_location = config['directory_id']
date = config['date']

# SAMPLES W/ PATHS TO LOOP THROUGH
samples = config['samples']

# Create a new presentation using premade template slide
prs = Presentation(config['powerpoint_template'])
print ('Creating new presentation')
runoff_data_path = None
runoff_dir = None
folder_path = None


# Unix & Windows paths are dif so determine what OS user is running 1st
if platform.system() == 'Windows':
    runoff_dir = r'\\nas01\Beadless_Results\gCAS4_Runoff_08132020\\'
    folder_path = r'\\nas01\Beadless_Results\gCAS4_Runoff_08132020\\'
    runoff_data_path = '\gCAS4_Runoff_08132020'
    print('You are on Windows!!!')
else:
    runoff_dir = '/Beadless_Results/gCAS4_Runoff_08132020'
    folder_path = '/Beadless_Results/gCAS4_Runoff_08132020/{}/gCAS4_Runoff_08132020'
    runoff_data_path = '/gCAS4_Runoff_08132020'
    print('You are on ', platform.system())

# Add images to slide accounting for both pre & post seq paths

def add_images(sample, placeholders, pic_path, report_path):
    # SNR PIC PLACEHOLDERS
    SNR_20 = None
    SNR_100 = None
    SNR_200 = None
    SNR_hist = None
    # SNR TEXT PLACEHOLDERS
    avg_snr = None
    snr_med = None
    snr_std = None
    med_jumps = None
    mean_noise = None

    if sample['pre_path']:
        SNR_20 = placeholders[18]
        SNR_100 = placeholders[19]
        SNR_200 = placeholders[20]
        SNR_hist = placeholders[21]
        
        avg_snr = placeholders[35]
        snr_med = placeholders[36]
        snr_std = placeholders[37]
        med_jumps = placeholders[38]
        mean_noise = placeholders[39]

    if sample['post_path']:
        SNR_20 = placeholders[14]
        SNR_100 = placeholders[15]
        SNR_200 = placeholders[16]
        SNR_hist = placeholders[17]

        avg_snr = placeholders[40]
        snr_med = placeholders[41]
        snr_std = placeholders[42]
        med_jumps = placeholders[43]
        mean_noise = placeholders[44]

    with open(report_path) as file:
        data = file.readlines()
        # print(round(float(data[20].split('=')[1][:-2]), 2))
        
        avg_snr.text = str(round(float(data[16].split('=')[1][:-2]), 2))
        snr_med.text = str(round(float(data[17].split('=')[1][:-2]), 2))
        snr_std.text = str(round(float(data[15].split('=')[1][:-2]), 2))
        med_jumps.text = str(round(float(data[19].split('=')[1][:-2]), 2))
        mean_noise.text = str(round(float(data[20].split('=')[1][:-2]), 2))

    # INSERTION PICS/DATA HERE
    SNR_20_pic = os.path.join(pic_path, 'SNR20_heatmap.png')
    SNR_100_pic = os.path.join(pic_path, 'SNR100_heatmap.png')
    SNR_200_pic = os.path.join(pic_path, 'SNR200_heatmap.png')
    SNR_hist_pic = os.path.join(pic_path, 'SNR_hist.png')

    SNR_20.insert_picture(SNR_20_pic)
    SNR_100.insert_picture(SNR_100_pic)
    SNR_200.insert_picture(SNR_200_pic)
    SNR_hist.insert_picture(SNR_hist_pic)

    loss_snr = placeholders[45]
    loss_jumps = placeholders[46]

    loss_snr.text = 'loss % SNR'
    loss_jumps.text = 'loss % jumps'


def loop_samples():
    for sample in samples:
        blank_slide_layout = prs.slide_masters[1].slide_layouts[9]
        # # Adds a slide using layout we picked
        slide = prs.slides.add_slide(blank_slide_layout)

        placeholders = slide.placeholders
        # ----------------- FXN TO GET PLACHOLDER IDX's which are not straightforward/easy to guess --------------------
        # for val in placeholders:
        #     print(val.name, val.placeholder_format.idx)
    
        pre_run_name = sample['pre_path']
        post_run_name = sample['post_path']
        pre_report_path = os.path.join(runoff_dir, pre_run_name, 'SNR report.txt')
        post_report_path = os.path.join(runoff_dir, post_run_name, 'SNR report.txt')

        pre_pic_path = '{}{}{}'.format(folder_path, pre_run_name, runoff_data_path)
        post_pic_path = '{}{}{}'.format(folder_path, post_run_name, runoff_data_path)


        if sample['pre_path']:
            add_images(sample, placeholders, pre_pic_path, pre_report_path)
        
        if sample['post_path']:
            add_images(sample, placeholders, post_pic_path, post_report_path)

        # Add General slide info
        id_placeholder = placeholders[26]
        label_placeholder = placeholders[27]
        pre_path_placeholder = placeholders[28]
        post_path_placeholder = placeholders[29]

        id_placeholder.text = sample['chip_id']
        label_placeholder.text = sample['chip_label']
        pre_path_placeholder.text = sample['pre_path']
        post_path_placeholder.text = sample['post_path']

def add_title_slide():
    slide = prs.slide_masters[1].slide_layouts[0]
    slide = prs.slides.add_slide(slide)
    placeholders = slide.placeholders

    print('Adding Title slide to presentation')

    # for val in placeholders:
    #     print(val.name, val.placeholder_format.idx)

    title = slide.shapes.title
    title.text = config['presentation_title']

    authors_placeholder = placeholders[13]
    date_placeholders = placeholders[14]

    authors_placeholder.text = config['authors']
    date_placeholders.text = 'Updated on {}'.format(config['updated_date'])


def add_conclusion_slide():
    slide = prs.slide_masters[1].slide_layouts[6]
    slide = prs.slides.add_slide(slide)
    placeholders = slide.placeholders
    print('Adding Conclusion slide to presentation')
    # for val in placeholders:
    #     print(val.name, val.placeholder_format.idx)

    conclusion_title = slide.shapes.title
    conclusion_title.text = 'Conclusion & next steps'


def upload_ppt(creds, file_name):
    drive_service = build('drive', 'v3', credentials=creds)

    # parents - ID of directory in Google Drive I want to upload file
    file_metadata = {'name': file_name, 'parents': [
        drive_save_location], 'mimetype': 'application/vnd.google-apps.presentation'}

    print('... Attempting to upload file to Google Drive')

    media = MediaFileUpload(
        file_name, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')

    file = drive_service.files().create(
        body=file_metadata, media_body=media, fields='id').execute()

    if file: 
        print('Presentation has been successfully uploaded to Google Drive')
        print('File ID: %s' % file.get('id'))


def auth_user():
    creds = None
    # The file token.pickle stores the user's access and refresh tokens, and is
    # created automatically when the authorization flow completes for the first
    # time.

    # If we already authenticated and token.pickle file already exists
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            # deserializing token from token.pickle with pickle module
            creds = pickle.load(token)
            # print('CREDS: ', pickle)
            print('User is authenticated')

    # If no token.pickle file or valid token, login and save credentials to a token.pickle file
    if not creds or not creds.valid:
        print('Attempting to authenticate user')
        if creds and creds.expired and creds.refresh_token:
            print('Getting new credentials for user')
            # manually refresh a credentials instance
            creds.refresh(Request())
        else:
            print('Attempting to create credentials for user')
            # Creates OAuth flow using credentials.json
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)
            # credentials2.json is Genapsys credentials
            # Opens browser to consent screen hosted on port 3000
            # If auth is valid, token.pickle is saved to computer
            # I think tokens are returned to us after browser authenticates
            creds = flow.run_local_server(port=3000)
        # Save the credentials for the next run
        with open('token.pickle', 'wb') as token:
            # dump() - Writes the pickled representation of the obj to the open file object file
            pickle.dump(creds, token)

    # FXN CALLS
    add_title_slide()
    loop_samples()
    add_conclusion_slide()
    file_name = 'runoff_{}.pptx'.format(date)
    prs.save(file_name)
    upload_ppt(creds, file_name)


auth_user()

# if __name__ == '__main__':
#     main()
